import os
import glob
import asyncio
import hypercorn.asyncio
import hypercorn.config
import re
import nltk
import weaviate
import psutil

from fastapi import FastAPI, HTTPException, Query
from fastapi.responses import StreamingResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from ollama import chat
from pypdf import PdfReader  # ✅ Use pypdf instead of PyPDF2
import docx
import pandas as pd
from sentence_transformers import SentenceTransformer
from transformers import pipeline
from rank_bm25 import BM25Okapi
from rdflib import Graph
from nltk.tokenize import sent_tokenize
from weaviate.connect import ConnectionParams
from weaviate.classes.config import Configure, Property
from weaviate.exceptions import WeaviateBaseError  # ✅ Correct Exception Handling
from sentence_transformers import SentenceTransformer  # ✅ Needed for encoding queries
from rank_bm25 import BM25Okapi  # ✅ Needed for BM25 scoring
from dotenv import load_dotenv
import tiktoken
import aiohttp
import asyncio
import json
import re
from pptx import Presentation  # Import pptx library

# Configuration
load_dotenv()
MY_FILES_DIR = "./my_files"
INSTRUCTIONS_FILE = "/app/instructions.txt"
SUPPORTED_FILE_TYPES = ["*.json", "*.pdf", "*.docx", "*.xlsx", "*.pptx", "*.ppt"]
LOCAL_MODEL_NAME = os.getenv("LOCAL_MODEL_NAME")
MODEL_NAME = os.getenv("MODEL_NAME")

# Get Ollama's external server details
TTYD_API_PORT = os.getenv("TTYD_API_PORT")
OLLAMA_TEMPERATURE = os.getenv("OLLAMA_TEMPERATURE")
OLLAMA_SERVER = f"http://ollama:11434"
OLLAMA_TEMPERATURE = os.getenv("OLLAMA_TEMPERATURE")
TTYD_UI_PORT = os.getenv("TTYD_UI_PORT")
LAST_N_CONVERSATION_TURNS = int(os.getenv("LAST_N_CONVERSATION_TURNS"))
TTYD_AGENTME_ENABLED = int(os.getenv("TTYD_AGENTME_ENABLED"))
AGENTME_API_URL = os.getenv("AGENTME_API_URL")
STRUCTURED_MODE = int(os.environ.get("STRUCTURED_RESPONSE_MODE"))

# Weaviate Configuration
WEAVIATE_HOST = "weaviate"
WEAVIATE_PORT = 8080
WEAVIATE_ALPHA = float(os.getenv("WEAVIATE_ALPHA"))  # ✅ Configurable hybrid search parameter

def calculate_dynamic_prompt_tokens():
    free_ram_gb = psutil.virtual_memory().available / (1024**3)
    print(f"Detected available system memory: {free_ram_gb:.2f} GB")

    HARDCAP_TOKENS = 8192
    RESERVED_COMPLETION_TOKENS = 512

    if free_ram_gb > 16:
        allowed_tokens = HARDCAP_TOKENS * 2
    else:
        allowed_tokens = HARDCAP_TOKENS

    allowed_tokens -= RESERVED_COMPLETION_TOKENS

    return allowed_tokens

MAX_PROMPT_TOKENS = calculate_dynamic_prompt_tokens()
print(f"Dynamic MAX_PROMPT_TOKENS = {MAX_PROMPT_TOKENS}")

# =============================================================================
# 2) MEMORY + SUMMARIES FOR LONG CONVERSATIONS
# =============================================================================

# >>> Session-based conversation logs <<<
session_logs: dict[str, list] = {}

def get_conversation_log(session_id: str) -> list:
    """
    Retrieve the conversation log for a given session_id.
    If it doesn't exist yet, initialize it to an empty list.
    """
    global session_logs
    if session_id not in session_logs:
        session_logs[session_id] = []
    return session_logs[session_id]

def clear_conversation_log(session_id: str):
    global session_logs
    session_logs[session_id] = []

import re
import aiohttp
import json

#
# Test Summarizer (restored to original prompt text)
#
async def summarize_text(text: str, temperature: float = 0.2) -> str:
    """
    Summarize large conversation text by calling your LLM (via Ollama).
    Removes references (e.g. [1], [2], etc.) from the final summary output.
    """
    prompt = f"Summarize or shorten the following text/question and do not involve in any discussion, just do it as better as possible and start immediately with summarized content, without introduction phrase:\n\n{text}"
    print("%%%%%%%")
    print(prompt)
    try:
        payload = {
            "model": MODEL_NAME,
            "prompt": prompt,
            "options": {
                "temperature": temperature,
                "num_ctx": 2048  # fewer tokens for summarizing
            },
            "stream": False
        }
        url = f"{OLLAMA_SERVER}/api/generate"

        async with aiohttp.ClientSession() as session:
            async with session.post(url, json=payload) as resp:
                result_json = await resp.json()
                # The final text is in result_json["response"]
                raw_summary = result_json.get("response", "").strip()

                # Remove bracketed references like [1], [2], etc. from the summary
                summary_no_refs = re.sub(r"\[\w+\]", "", raw_summary)

                return summary_no_refs
    except Exception as e:
        print(f"Error summarizing text: {e}")
        return "Summary unavailable (error)."


# =============================================================================
# 3) FASTAPI SETUP
# =============================================================================

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"], 
    allow_headers=["*"], 
)

def ensure_collection_exists(client, collection_name="Document"):
    try:
        return client.collections.get(collection_name)
    except WeaviateBaseError:  # ✅ Catch exception when collection does not exist
        return client.collections.create(
            name=collection_name,
            configure_vectorizer=True,  # ✅ Required for text2vec-transformers
            vectorizer_config=Configure.Vectorizer.text2vec_transformers(
                model=LOCAL_MODEL_NAME,  # ✅ Make sure it matches your SentenceTransformer model
                pooling_strategy="cls",
                vectorize_property_name=True  # ✅ Enable automatic vectorization of `content`
            ),
            properties=[
                Property(name="content", data_type="text"),
                Property(name="source", data_type="text"),
            ]
        )


# ✅ Set Weaviate connection details (local-only, no external API calls)
client = weaviate.WeaviateClient(
    connection_params=ConnectionParams.from_params(
        http_host=WEAVIATE_HOST,
        http_port=int(WEAVIATE_PORT),
        http_secure=False,
        grpc_host=WEAVIATE_HOST,
        grpc_port=50051,
        grpc_secure=False,
    ),
    skip_init_checks=False
)

client.connect()

# ✅ Load instructions
def load_instructions():
    try:
        with open(INSTRUCTIONS_FILE, "r", encoding="utf-8") as f:
            return f.read().strip()
    except Exception as e:
        print(f"⚠️ Warning: Could not load instructions.txt ({e}). Using default instructions.")
        return "You are an AI assistant. Answer questions based on the given content."

INSTRUCTIONS = load_instructions()

# ✅ Chunking function
def chunk_text(text, max_chunk_size=512):
    sentences = sent_tokenize(text)
    chunks, current_chunk, current_length = [], [], 0

    for sentence in sentences:
        if current_length + len(sentence) > max_chunk_size:
            chunks.append(" ".join(current_chunk))
            current_chunk, current_length = [], 0
        current_chunk.append(sentence)
        current_length += len(sentence)
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    return chunks

# ✅ Read supported files
def read_supported_files(directory):
    documents = []
    for file_type in SUPPORTED_FILE_TYPES:
        for filepath in glob.glob(os.path.join(directory, file_type)):
            # Skip temporary files that start with '~$'
            if os.path.basename(filepath).startswith("~$"):
                print(f"skip temporary file: {filepath}")
                continue

            content = None  # initialize content to None

            if filepath.endswith(".json"):
                with open(filepath, "r", encoding="utf-8", errors="replace") as file:
                    content = file.read()
            elif filepath.endswith(".pdf"):
                content = "\n".join([page.extract_text() or "" for page in PdfReader(filepath).pages])
            elif filepath.endswith(".docx"):
                content = "\n".join([para.text for para in docx.Document(filepath).paragraphs])
            elif filepath.endswith(".xlsx"):
                excel_data = pd.ExcelFile(filepath)
                content = "\n".join(excel_data.parse(sheet).to_string(index=False, header=True) for sheet in excel_data.sheet_names)
            elif filepath.endswith(".pptx") or filepath.endswith(".ppt"):
                # Handle PPTX and PPT files
                presentation = Presentation(filepath)
                
                # Iterate over slides
                for slide_idx, slide in enumerate(presentation.slides):
                    content = ""
                    
                    # Iterate over shapes to collect text from all text-containing shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):  # Check if the shape contains text
                            # Clean up extra new lines and spaces
                            cleaned_text = shape.text.strip().replace('\r', '')  # Remove carriage returns
                            cleaned_text = re.sub(r'(\n\s*){2,}', '\n', cleaned_text)  # Replace multiple new lines with a single new line
                            content += cleaned_text + " "  # Add space between texts from shapes
                    
                    # Instead of chunking the text, store the whole slide text as one chunk
                    documents.append({"content": content.strip(), "source": f"{os.path.basename(filepath)}_slide_{slide_idx + 1}_chunk_0"})
                    print("------------")
                    print(content)
                    print("------------------")
                # Clear content to avoid re-processing the last slide's text
                content = None
            else:
                content = None

            if content:
                if filepath.endswith(".json"):
                    # For JSON files, use the full content without chunking.
                    documents.append({"content": content, "source": f"{os.path.basename(filepath)}_chunk_0"})
                else:
                    # For other file types, chunk the content.
                    for idx, chunk in enumerate(chunk_text(content)):
                        print("-----------")
                        print(chunk)
                        print("-----------")
                        documents.append({"content": chunk, "source": f"{os.path.basename(filepath)}_chunk_{idx}"})
    return documents

# ✅ Embed and store in Weaviate
def embed_and_store(documents, client, collection_name="Document"):
    embedding_model = SentenceTransformer(LOCAL_MODEL_NAME)
    collection = ensure_collection_exists(client, collection_name)
    for doc in documents:
        embedding = embedding_model.encode(doc["content"]).tolist()
        collection.data.insert(
            properties={
                "content": doc["content"],
                "source": doc["source"],
                "vector": embedding
            }
        )

# ✅ BM25 Indexing
def tokenize(text):
    return text.lower().split()

def build_bm25_index(documents):
    return BM25Okapi([tokenize(doc["content"]) for doc in documents]), documents



# Initialize Weaviate and document indexing
print("Initializing Weaviate client...")
documents = read_supported_files(MY_FILES_DIR)
print(f"Embedding {len(documents)} document chunks into Weaviate...")
embed_and_store(documents, client)
bm25, indexed_docs = build_bm25_index(documents)  # ✅ Ensure these variables are set


# ✅ Hybrid retrieval
def hybrid_retrieval(question, top_k=10):
    embedding_model = SentenceTransformer(LOCAL_MODEL_NAME)
    question_embedding = embedding_model.encode(question).tolist()

    collection = client.collections.get("Document")
    response = collection.query.hybrid(
        query=question,
        vector=question_embedding,
        alpha=WEAVIATE_ALPHA,
        limit=top_k,
        return_properties=["content", "source"]
    )

    vector_chunks = response.objects
    if not vector_chunks:
        return [], "No relevant information found."

    retrieved_docs = [
        {"content": doc.properties["content"], "source": doc.properties["source"]}
        for doc in vector_chunks
    ]
    return retrieved_docs, None

class QuestionRequest(BaseModel):
    # Session ID is optional from the client; default to 0 if missing
    session_id: str = "0"
    question: str

def build_full_prompt(session_id: str, question: str, retrieved_docs: list) -> str:
    # Use the conversation log for this session
    conversation_log = get_conversation_log(session_id)

    # Deduplicate retrieved content by file
    unique_content = {}
    for doc in retrieved_docs:
        file_name = doc["source"].split("_chunk_")[0]
        if file_name not in unique_content:
            unique_content[file_name] = doc["content"]

    formatted_context = "\n\n".join(
        [f"[{i+1}] {txt} (Source: {fn})" for i, (fn, txt) in enumerate(unique_content.items())]
    )

    # Gather the last N turns from conversation_log
    last_n_interactions = conversation_log[-(2*LAST_N_CONVERSATION_TURNS):]

    conversation_text = ""
    for item in last_n_interactions:
        role = item["role"]
        content = item["content"]
        conversation_text += f"{role.capitalize()}: {content}\n"

    final_prompt = f"""
{INSTRUCTIONS}

# Recent Conversation
{conversation_text}

# Retrieved Documents
{formatted_context}

User's new question: {question}

Answer:
Please generate an answer based on all relevant context. Cite sources using [number] brackets based on the sources in the retrieved documents.
"""

    return final_prompt.strip()

async def ensure_prompt_within_limit(session_id: str, question: str, retrieved_docs: list) -> str:
    conversation_log = get_conversation_log(session_id)
    draft_prompt = build_full_prompt(session_id, question, retrieved_docs)
    encoding = tiktoken.get_encoding("cl100k_base")
    token_count = len(encoding.encode(draft_prompt))

    if token_count <= MAX_PROMPT_TOKENS:
        return draft_prompt

    print(f"⚠️ Prompt length {token_count} exceeds limit {MAX_PROMPT_TOKENS}. Summarizing older conversation...")

    if len(conversation_log) <= 1:
        return ""

    # Summarize everything up to the last user question
    old_part = conversation_log[:-1]
    last_message = conversation_log[-1]  # the user’s current question

    # Use the local summarizer
    old_text = "\n".join([f'{x["role"].capitalize()}: {x["content"]}' for x in old_part])
    summary = await summarize_text(old_text)

    # Keep only the last user question in the log
    clear_conversation_log(session_id)
    get_conversation_log(session_id).append(last_message)

    new_draft = build_full_prompt(session_id, question, retrieved_docs)
    new_count = len(encoding.encode(new_draft))

    if new_count > MAX_PROMPT_TOKENS:
        print(f"⚠️ Even after summarizing, prompt length {new_count} > {MAX_PROMPT_TOKENS}.")
        return ""

    return new_draft

def collect_source_files(retrieved_docs: list) -> dict:
    """
    Build a dictionary {file_name: content} matching the same indexing
    done in build_full_prompt(). The index i+1 will correspond to [i+1].
    """
    unique_content = {}
    for doc in retrieved_docs:
        file_name = doc["source"].split("_chunk_")[0]
        if file_name not in unique_content:
            unique_content[file_name] = doc["content"]
    return unique_content

# ---------------- NEW: Call AgentMe if enabled ----------------
async def call_agentme_api(question: str) -> str:
    """
    Call AgentMe's /chat endpoint with the user's question,
    return the text of the response or a fallback if error.
    """
    payload = {"message": question}
    try:
        async with aiohttp.ClientSession() as session:
            async with session.post(f"{AGENTME_API_URL}/chat", json=payload) as resp:
                if resp.status == 200:
                    data = await resp.json()
                    # Expecting: {"response":"some text"}
                    return data.get("response", "")
                else:
                    return f"[AgentMe error: status {resp.status}]"
    except Exception as e:
        return f"[Error calling AgentMe: {e}]"

async def generate_answer_with_ollama(session_id: str, question: str):
    # 1) Retrieve the conversation log for this session
    conversation_log = get_conversation_log(session_id)

    # 2) Append the user message to that log
    conversation_log.append({"role": "user", "content": question})

    # 3) Hybrid retrieval
    retrieval_query = (
        f"Recent messages:\n{question}\n"
        f"Now answer this question: {question}"
    )
    retrieved_docs, retrieval_error = hybrid_retrieval(retrieval_query)
    if retrieval_error:
        if STRUCTURED_MODE == 1:
            # Return JSON with an error as the LLM response
            yield json.dumps({
                "llm_response": retrieval_error,
                "action": ""
            })
        else:
            # Return the error as plain text
            yield retrieval_error
        return

    # 4) Build the final prompt
    final_prompt = await ensure_prompt_within_limit(session_id, question, retrieved_docs)
    if not final_prompt:
        if STRUCTURED_MODE == 1:
            # Return JSON with the error
            yield json.dumps({
                "llm_response": "I'm sorry, but the conversation is too large to process. Please start a new session.",
                "action": ""
            })
        else:
            yield "I'm sorry, but the conversation is too large to process. Please start a new session."
        return

    # Collect the source mapping so we can display them at the end
    unique_content = collect_source_files(retrieved_docs)

    print("##############################################################################################")
    print("--- FINAL PROMPT ---")
    print(final_prompt)
    print("--------------------")

    # We will collect the full LLM response here no matter what
    assistant_full_response = ""

    try:
        encoding = tiktoken.get_encoding("cl100k_base")
        prompt_tokens = len(encoding.encode(final_prompt))
        print(f"Prompt tokens: {prompt_tokens} / limit {MAX_PROMPT_TOKENS}")

        payload = {
            "model": MODEL_NAME,
            "prompt": final_prompt,
            "options": {
                "temperature": float(OLLAMA_TEMPERATURE),
                "num_ctx": MAX_PROMPT_TOKENS
            },
            "stream": True
        }

        url = f"{OLLAMA_SERVER}/api/generate"

        async with aiohttp.ClientSession() as session:
            async with session.post(url, json=payload) as resp:
                if STRUCTURED_MODE == 1:
                    # Accumulate the entire response without yielding partial chunks
                    async for line in resp.content:
                        line = line.decode('utf-8').strip()
                        if not line:
                            continue
                        try:
                            data = json.loads(line)
                            chunk_text = data.get("response", "")
                            if chunk_text:
                                assistant_full_response += chunk_text
                        except json.JSONDecodeError:
                            print("Error decoding JSON:", line)
                else:
                    # Stream partial chunks to the user in real-time
                    buffer = ""
                    async for line in resp.content:
                        line = line.decode('utf-8').strip()
                        if not line:
                            continue
                        try:
                            data = json.loads(line)
                            chunk_text = data.get("response", "")
                            if chunk_text:
                                assistant_full_response += chunk_text
                                buffer += chunk_text

                                words = re.findall(r'\S+\s*', buffer)
                                # Yield every complete word except the last partial one
                                for word in words[:-1]:
                                    yield word
                                    await asyncio.sleep(0.01)
                                buffer = words[-1] if words else ""
                        except json.JSONDecodeError:
                            print("Error decoding JSON:", line)
                    # If there's leftover buffer text, yield that as the final chunk
                    if buffer:
                        yield buffer

        # 5) If AgentMe is enabled, call it for the same question
        agent_me_resp = ""
        if TTYD_AGENTME_ENABLED == 1:
            last_n_interactions = conversation_log[-(2*LAST_N_CONVERSATION_TURNS):]
            conversation_text = ""
            for item in last_n_interactions:
                role = item["role"]
                content = item["content"]
                conversation_text += f"{role.capitalize()}: {content}\n"
            agent_me_resp = await call_agentme_api(conversation_text)

        # 6) Summarize the final response for storing in conversation history
        summarized_answer = await summarize_text(assistant_full_response)
        conversation_log.append({"role": "assistant", "content": summarized_answer})

        # 7) Build sources text
        sources_text = "\n\nSources:\n"
        for i, (fn, _) in enumerate(unique_content.items()):
            sources_text += f"[{i+1}] {fn}\n"

        # 8) Final output
        if STRUCTURED_MODE == 1:
            # In structured mode, produce a single JSON with two fields:
            #   1) llm_response (the full text from the LLM)
            #   2) action (the AgentMe response, if any)
            # Optionally, you could store the sources in either field
            structured_json = {
                "llm_response": assistant_full_response + sources_text,
                "action": agent_me_resp
            }
            yield json.dumps(structured_json)
        else:
            # Existing approach: Output final AgentMe text, plus sources
            if TTYD_AGENTME_ENABLED == 1 and agent_me_resp:
                yield f"\n\n--- AgentMe response ---\n{agent_me_resp}\n"
            yield sources_text

    except Exception as e:
        if STRUCTURED_MODE == 1:
            yield json.dumps({
                "llm_response": f"An error occurred: {e}",
                "action": ""
            })
        else:
            yield f"An error occurred: {e}"
    finally:
        print("##############################################################################################")

# Health check endpoint
@app.get("/health")
async def health_check():
    return JSONResponse(content={"status": "ok"})

# >>> CHANGED TO GET + optional session_id <<<
@app.get("/clear_session")
async def clear_session(session_id: str = Query("0")):
    """
    Clears the conversation history for the specified session ID (default=0).
    """
    clear_conversation_log(session_id)
    return JSONResponse(content={"message": f"Conversation log cleared for session_id={session_id}."})

@app.post("/ask")
async def ask(request: QuestionRequest):
    """
    The user sends { "question": "...", "session_id": "123" }.
    If session_id is omitted, we treat it as 0.
    """
    if not request.question:
        raise HTTPException(status_code=400, detail="Question cannot be empty")
    
    # Extract session_id (or default to "0")
    session_id = request.session_id or "0"

    # Pass it into the streaming generator
    return StreamingResponse(generate_answer_with_ollama(session_id, request.question), media_type="text/plain")


# Determine the number of CPU cores and set the number of workers accordingly
cpu_cores = os.cpu_count()  # Get the number of CPU cores
workers = cpu_cores if cpu_cores is not None else 1  # Use number of cores, default to 1 if unavailable

# Set Hypercorn configuration with workers based on the number of CPU cores
config = hypercorn.config.Config()
config.bind = [f"0.0.0.0:{TTYD_API_PORT}"]  # Ensure binding is correct
config.alpn_protocols = ["h2", "http/1.1"]  # Explicitly allow HTTP/2 first
config.workers = workers  # Set the number of workers based on the number of CPU cores
config.accesslog = "-"

# Run the Hypercorn server with multiprocessing support (using asyncio)
if __name__ == "__main__":
    print(f"🚀 Starting server with {workers} workers on port {TTYD_API_PORT} ...")
    # Use the correct asyncio method to start the server
    asyncio.run(hypercorn.asyncio.serve(app, config))
